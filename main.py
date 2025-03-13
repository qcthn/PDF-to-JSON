import streamlit as st
import pdfplumber
import json
import os
import shutil
import openai
from docx import Document
import re
import pandas as pd
from pptx import Presentation
LOGO_URL_LARGE = "logo-medium.png"
st.image(
    # LOGO_URL_LARGE,
    # size = "large",
    # link="https://streamlit.io/gallery",
    LOGO_URL_LARGE,
)
# Nhập OpenAI API Key
# api_key = st.text_input("🔑 Enter OpenAI API Key:", type="password")

# if not api_key:
#     st.warning("Please enter OpenAI API Key to use the app.")
#     st.stop()
if "openai" in st.secrets and "api_key" in st.secrets["openai"]:
    api_key = st.secrets["openai"]["api_key"]
else:
    api_key = st.text_input("Nhập OpenAI API Key:", type="password")

if api_key:
    st.write("✅ API Key đã được nhập!")

# Cấu hình OpenAI client
client = openai.OpenAI(api_key=api_key)
def clean_text(text):
    """Làm sạch văn bản để loại bỏ ký tự NULL và các ký tự điều khiển không hợp lệ."""
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)  # Loại bỏ ký tự điều khiển
    return text.strip()
# Làm sạch JSON phản hồi từ OpenAI
def clean_json_response(response_text):
    """Làm sạch phản hồi GPT để loại bỏ các ký tự không mong muốn trước khi phân tích JSON."""
    response_text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', response_text)
    response_text = response_text.strip()
    if response_text.startswith("```json"):
        response_text = response_text[7:]
    if response_text.endswith("```"):
        response_text = response_text[:-3]
    return response_text.strip()
# Hàm trích xuất văn bản từ PDF
def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += clean_text(page_text) + "\n"
    except Exception as e:
        st.error(f"Lỗi khi xử lý PDF: {e}")
        return None
    return text
# Trích xuất văn bản + bảng từ DOCX
def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        full_text = []

        # Trích xuất văn bản từ đoạn (paragraphs)
        for paragraph in doc.paragraphs:
            full_text.append(clean_text(paragraph.text))

        # Trích xuất dữ liệu từ bảng
        for table in doc.tables:
            for row in table.rows:
                row_text = [clean_text(cell.text) for cell in row.cells]
                full_text.append(" | ".join(row_text))  # Ngăn cách bằng " | " để giữ định dạng

        return "\n".join(full_text)
    except Exception as e:
        st.error(f"Error processing DOCX: {e}")
        return None

# Trích xuất văn bản từ TXT
def extract_text_from_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return clean_text(f.read())
    except Exception as e:
        st.error(f"Error processing TXT: {e}")
        return None

# Trích xuất dữ liệu từ EXCEL (giữ đúng thứ tự cột + hàng)
def extract_data_from_excel(file_path):
    try:
        df = pd.read_excel(file_path, dtype=str)  # Đọc tất cả dữ liệu dưới dạng chuỗi
        text_output = df.to_string(index=False)  # Chuyển DataFrame thành văn bản dễ đọc
        return text_output
    except Exception as e:
        st.error(f"Error processing Excel: {e}")
        return None
# Trích xuất dữ liệu từ PPT
def extract_text_from_pptx(file_path):
    """
    Trích xuất toàn bộ văn bản từ file PPTX.
    Gồm cả văn bản trong Slide, Shape và Table.
    """
    try:
        prs = Presentation(file_path)
        text_content = []

        for slide in prs.slides:
            # Lặp qua tất cả shape
            for shape in slide.shapes:
                # Nếu shape có text_frame (vd: text box), thì đọc text
                if shape.has_text_frame:
                    text_content.append(clean_text(shape.text))

                # Nếu shape chứa bảng, ta duyệt qua bảng đó để đọc text
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            row_text.append(clean_text(cell.text))
                        text_content.append(" | ".join(row_text))

        return "\n".join(text_content)
    except Exception as e:
        st.error(f"Error processing PPTX: {e}")
        return None
# Xử lý đầu vào nhiều loại tệp
def extract_text_from_file(file_path, file_type):
    if file_type == "pdf":
        return extract_text_from_pdf(file_path)
    elif file_type == "docx":
        return extract_text_from_docx(file_path)
    elif file_type == "txt":
        return extract_text_from_txt(file_path)
    elif file_type == "xlsx":
        return extract_data_from_excel(file_path)
    elif file_type == "pptx":  
        return extract_text_from_pptx(file_path)
    else:
        st.error("Unsupported file format.")
        return None

# Hàm gọi GPT để trích xuất thông tin
def extract_info_with_gpt(text):
    prompt = f"""
    Trích xuất thông tin từ văn bản CV và trả về JSON hợp lệ:
    {{
        "Name": "",
        "Email": "",
        "Phone": "",
        "Skills": [],  
        "Experience": [],  
        "Education": [],  
        "Certifications": [],  
        "Languages": [],  
        "Strengths": [],  
        "Weaknesses": [],  
        "Additional information": []
    }}
    For the **Languages** field, include:
    - The candidate's native language based on their nationality (e.g., Vietnamese for a candidate from Vietnam).
    - Any foreign language certifications (e.g., TOEIC score) and the corresponding language proficiency level (e.g., English with a proficiency level based on the score).

    For **Strengths and Weaknesses**, analyze the candidate's work experience to identify:
    - **Strengths:** Key skills and attributes demonstrated through their experience.
    - **Weaknesses:** Areas for improvement or challenges faced in their roles.
    CV text:
    {text}
    """

    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are an expert in extracting information from CVs (resumes) and images with 10 years of experience in getting the exact information needed to recruit suitable positions for the company."

            "**Context:** I will provide you with resumes of candidates (which can be one or more) or image files containing text."

            "**Your task** is to extract information from the resumes and images I provide (I have taken the text from the resume, and the image will be provided to you below) and return the output as a JSON file."

            "Some of the most important information required for each candidate includes:"
            "- Name"
            "- Email"
            "- Phone number"
            "- Skills"
            "- Experience (including: position, timeline, responsibilities)"
            "- Education (including: degree, institution, timeline, GPA)"
            "- Certifications"
            "- Languages (including proficiency based on nationality and language certifications)"
            "- Strengths (based on the candidate's experience and job description)"
            "- Weaknesses (based on the candidate's experience and job description)"
            "- Additional information (including identification and visa details if provided)"

            "**Task:** Extract the following information from the CV text and return it as JSON."

            "**Output:** JSON file format"

            "***Note:** I can provide you with the text, but in that text will be a synthesis of many resumes of different candidates.*"

            "**REMEMBER:** The output should only be in JSON format." },
            {"role": "user", "content": prompt}
        ]
    )
    # return json.loads(response.choices[0].message.content.strip())
    extracted_text = response.choices[0].message.content.strip()
    cleaned_text = clean_json_response(extracted_text)
    # return json.loads(response.choices[0].message.content.strip())
    return json.loads(cleaned_text)
# Lưu JSON
def save_to_json(data_list, output_path):
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data_list, f, ensure_ascii=False, indent=4)
@st.cache_data
def generate_json():
    extracted_data = {}
    for filename, text in extracted_texts.items():
        extracted_info = extract_info_with_gpt(text)  # Gọi GPT tự động
        extracted_info["Filename"] = filename
        extracted_data[filename] = extracted_info

    return json.dumps(list(extracted_data.values()), ensure_ascii=False, indent=4)
# Tạo file Word
def create_word_file(text_list, output_path):
    try:
        doc = Document()
        for idx, (filename, text) in enumerate(text_list):
            doc.add_heading(f"File {idx + 1}: {filename}", level=1)
            doc.add_paragraph(clean_text(text))
            doc.add_page_break()
        doc.save(output_path)
        return True
    except Exception as e:
        st.error(f"Error creating Word file: {e}")
        return False

# Ứng dụng Streamlit
st.title("📄 HRIS AI Assistant")
st.write("🔹 Upload your CV (PDF), extract content and ask questions to the virtual assistant!")
    # Tạo thư mục tạm
temp_dir = "temp"
os.makedirs(temp_dir, exist_ok=True)

uploaded_files = st.file_uploader("Upload PDFs", accept_multiple_files=True)

extracted_texts = {}
extracted_data = {}
if uploaded_files:
    for uploaded_file in uploaded_files:
        temp_file_path = os.path.join(temp_dir, uploaded_file.name)
        with open(temp_file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        file_type = uploaded_file.name.split(".")[-1].lower()
        text = extract_text_from_file(temp_file_path, file_type)
        # text = extract_text_from_pdf(temp_file_path)
        if text:
            extracted_texts[uploaded_file.name] = text  # Lưu văn bản chưa xử lý
        os.remove(temp_file_path)

    # Hiển thị nội dung PDF
    for filename, text in extracted_texts.items():
        st.text_area(f"Content ({filename})", text)

    # Khi người dùng ấn nút, thực hiện extract và tạo JSON
    if st.button("📥 Generate & Download JSON file"):
        for filename, text in extracted_texts.items():
            extracted_info = extract_info_with_gpt(text)  # Chỉ gọi GPT khi cần
            extracted_info["Filename"] = filename
            extracted_data[filename] = extracted_info

        json_output = os.path.join(temp_dir, "extracted_data.json")
        json_data = json.dumps(list(extracted_data.values()), ensure_ascii=False, indent=4)

        # Tạo link tải xuống ngay lập tức
        st.download_button(
            label="📥 Click here to download JSON file",
            data=json_data,
            file_name="extracted_data.json",
            mime="application/json"
        )
# Điểm khác biệt quan trọng:
 

 
    # Chatbot với trợ lý ảo
    st.subheader("💬  Chat with virtual assistant")

    if len(uploaded_files) > 1:
        selected_cv = st.selectbox("Select CV to interact:", list(extracted_texts.keys()))
    else:
        selected_cv = list(extracted_texts.keys())[0]
    if "openai_model" not in st.session_state:
        st.session_state.openai_model = "gpt-3.5-turbo"
    if "messages" not in st.session_state:
        st.session_state.messages = []
    # Display chat messages from history on app rerun
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    if "openai_model" not in st.session_state:
        st.session_state["openai_model"] = "gpt-3.5-turbo"
    text_CV = f"Below is the content of the candidate's CV:\n{extracted_texts[selected_cv]}"
    prompt = st.chat_input("Ask a virtual assistant about this resume:")
    if prompt:
         # Add user message to chat history
        st.session_state.messages.append({"role": "user", "content": prompt})
        # Display user message in chat message container
        with st.chat_message("user"):
            st.markdown(prompt)
        with st.chat_message("assistant"):
            stream = client.chat.completions.create(
            model=st.session_state["openai_model"],
            messages=[{"role": "system", "content": "Play the role of a professional HR, with 10 years of experience in finding potential candidates suitable for the company based on the CV (resume) they send Context: I will provide you with information of each CV (resume) in text form, from which I will ask you some questions related to the CV (resume) of this candidate Task: Please provide the most accurate and closest information to the question I asked, helping me have the most objective view of this candidate so that I can decide whether to hire him or not Tone: solemn, dignified, straightforward, suitable for the office environment, recruitment. Below is the content of the candidate's CV"  + extracted_texts[selected_cv] },
                {"role": "user", "content": prompt}],
                stream=True,
            )
                # if hasattr(response.choices[0].delta, "content"):    
            response = st.write_stream(stream)
        st.session_state.messages.append({"role": "assistant", "content": response})
    shutil.rmtree(temp_dir, ignore_errors=True)
